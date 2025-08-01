import streamlit as st
import google.generativeai as genai
from docx import Document
import pdfplumber
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
import os
from dotenv import load_dotenv

# ---------- Load API Key ----------
load_dotenv()
api_key = os.getenv("GEMINI_API_KEY") or st.secrets["gemini"]["api_key"]
genai.configure(api_key=api_key)
model = genai.GenerativeModel("gemini-2.5-flash")

# ---------- Document Readers ----------
def read_docx(file):
    doc = Document(file)
    return "\n".join([para.text.strip() for para in doc.paragraphs if para.text.strip()])

def read_pdf(file):
    with pdfplumber.open(file) as pdf:
        return "\n".join([page.extract_text() for page in pdf.pages if page.extract_text()])

def read_pptx(file):
    prs = Presentation(file)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text.strip())
    return "\n".join(texts)

def read_txt(file):
    return file.read().decode("utf-8")

def scrape_website(url):
    try:
        response = requests.get(url, timeout=10)
        soup = BeautifulSoup(response.text, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)
    except Exception as e:
        return f"Error scraping website: {str(e)}"

# ---------- Prompt Templates ----------
def load_sales_prompt(platform):
    default_prompts = {
        "LinkedIn": """🔗 LinkedIn Message Prompt:
You are a senior B2B sales strategist. Write a short, human, personalized LinkedIn message. Make it relevant, clear, and rooted in uploaded product content. No buzzwords or filler. Max 70 words.""",

        "WhatsApp": """💬 WhatsApp Message Prompt:
Write a casual yet professional WhatsApp message introducing a B2B SaaS product. Keep it brief (max 50 words), friendly, and based only on uploaded content. Don’t use emojis or sales jargon.""",

        "Email": """✉️ Email Prompt:
Craft a short outbound B2B email introducing a product. Include a strong opener, clear value prop, 2–3 features, and a soft CTA. Stay under 140 words. Use uploaded materials only."""
    }

    filename_map = {
        "LinkedIn": "linkdln.docx",
        "WhatsApp": "whatsapp.docx",
        "Email": "email.docx"
    }
    file_path = filename_map.get(platform, "")
    if file_path and os.path.exists(file_path):
        return read_docx(file_path)

    return default_prompts.get(platform, "Message format not found.")

# ---------- Streamlit UI ----------
st.set_page_config(page_title="AI Message Generator", layout="wide")
st.title("🧠 AI Email, LinkedIn & WhatsApp Message Generator")

platform = st.selectbox("Choose Message Platform", ["LinkedIn", "WhatsApp", "Email"])
message_type = st.selectbox("Select Message Type", ["Generic", "Personalized"])

col1, col2 = st.columns(2)

user_content = ""
custom_prompt = ""
generated_message = ""

# ---------- Left Column: Upload & Prompt ----------
with col1:
    st.subheader("📁 Upload a Document")
    uploaded_file = st.file_uploader("Upload a document", type=["docx", "pdf", "pptx", "txt"])
    if uploaded_file:
        ext = uploaded_file.name.split(".")[-1]
        if ext == "docx":
            user_content = read_docx(uploaded_file)
        elif ext == "pdf":
            user_content = read_pdf(uploaded_file)
        elif ext == "pptx":
            user_content = read_pptx(uploaded_file)
        elif ext == "txt":
            user_content = read_txt(uploaded_file)
        else:
            st.error("Unsupported file type.")

    if message_type == "Personalized":
        st.subheader("✍️ Write Your Prompt")
        custom_prompt = st.text_area(
            f"Describe the kind of {platform} message you want (tone, purpose, etc.):",
            placeholder=f"e.g., Write a concise, warm {platform} message to a CMO about our AI product that improves customer engagement."
        )

# ---------- Right Column: Personalization ----------
with col2:
    if message_type == "Personalized":
        st.subheader("👤 Personalization Fields")
        recipient_name = st.text_input("Recipient Name")
        company_name = st.text_input("Company Name")
        job_title = st.text_input("Job Title (Optional)")

# ---------- Generate Button ----------
if st.button(f"🧠 Generate {platform} Message"):
    if not user_content.strip() and not custom_prompt.strip():
        st.error("❗ Please upload a document or enter a custom prompt before generating the message.")
    elif message_type == "Personalized" and (not recipient_name.strip() or not company_name.strip()):
        st.error("❗ Please fill in recipient name and company name for a personalized message.")
    else:
        with st.spinner("Generating..."):
            prompt_template = load_sales_prompt(platform)

            if message_type == "Generic":
                final_prompt = (
                    f"{prompt_template}\n\n"
                    f"Generate a **generic {platform} message** using this content:\n\n{user_content[:10000]}"
                )
            elif custom_prompt:
                final_prompt = (
                    f"{custom_prompt}\n\n(Here is some background info to help craft the {platform} message:)\n\n{user_content[:10000]}"
                )
            else:
                final_prompt = (
                    f"{prompt_template}\n\n"
                    f"Generate a **personalized {platform} message** for:\n"
                    f"Name: {recipient_name}\n"
                    f"Company: {company_name}\n"
                    f"Job Title: {job_title}\n\n"
                    f"Background:\n{user_content[:10000]}"
                )

            try:
                response = model.generate_content(final_prompt)
                st.session_state.generated_message = response.text
                st.session_state.edited_message = ""  # Reset edit
                st.session_state.show_edit = True
                st.success(f"✅ {platform} Message Generated!")
            except Exception as e:
                st.error(f"Generation failed: {e}")
                st.session_state.generated_message = ""

# ---------- Show Generated Message ----------
if st.session_state.get("generated_message"):
    st.subheader(f"📨 Your {platform} Message")
    st.write(st.session_state.generated_message)

    # ---------- Post-Generation Editing ----------
    if st.session_state.get("show_edit"):
        st.subheader("✏️ Refine the Message")
        edit_instruction = st.text_area(
            "Want to tweak it? Type your edit instructions here:",
            placeholder="e.g., Make it shorter and more casual"
        )

        if st.button("✂️ Edit Message"):
            if edit_instruction.strip():
                with st.spinner("Editing message..."):
                    edit_prompt = (
                        f"Please revise the following {platform} message based on this instruction:\n"
                        f"{edit_instruction.strip()}\n\n"
                        f"Original Message:\n{st.session_state.generated_message}"
                    )
                    try:
                        edit_response = model.generate_content(edit_prompt)
                        st.session_state.edited_message = edit_response.text
                        st.success("✅ Edited Message Ready")
                    except Exception as e:
                        st.error(f"Editing failed: {e}")
            else:
                st.warning("Please provide some instructions before editing.")

# ---------- Show Edited Message ----------
if st.session_state.get("edited_message"):
    st.subheader("✅ Refined Message")
    st.write(st.session_state.edited_message)
